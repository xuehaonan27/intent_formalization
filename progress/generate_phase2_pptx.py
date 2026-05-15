#!/usr/bin/env python3
"""Generate the spec-determinism Phase 2 progress deck.

The visual style intentionally follows progress/spec-testing-update.pptx:
white background, teal section titles, thin divider line, and concrete
example slides with two code panels.

Dependency:
    python3 -m pip install python-pptx

Usage:
    python3 progress/generate_phase2_pptx.py
    python3 progress/generate_phase2_pptx.py --output /tmp/phase2.pptx
"""

from __future__ import annotations

import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ModuleNotFoundError as exc:
    raise SystemExit(
        "Missing dependency: python-pptx\n"
        "Install it with: python3 -m pip install python-pptx"
    ) from exc


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TEAL = RGBColor(0, 121, 107)
DARK = RGBColor(45, 45, 45)
MUTED = RGBColor(95, 95, 95)
CODE_BG = RGBColor(241, 241, 246)
CODE_BORDER = RGBColor(190, 190, 198)
SHADOW = RGBColor(218, 218, 224)
WHITE = RGBColor(255, 255, 255)
PALE_TEAL = RGBColor(232, 247, 244)
RED = RGBColor(176, 57, 57)

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"
CODE_FONT = "Courier New"


def add_text(slide, x, y, w, h, text, *, size=20, bold=False, italic=False,
             color=DARK, font=BODY_FONT, align=None):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return shape


def add_header(slide, title: str) -> None:
    add_text(slide, Inches(0.62), Inches(0.32), Inches(12.0), Inches(0.55),
             title, size=29, bold=True, color=TEAL, font=TITLE_FONT)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.52), Inches(0.95), Inches(12.1), Inches(0.025)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def add_footer(slide, number: int) -> None:
    add_text(slide, Inches(0.65), Inches(7.1), Inches(5.0), Inches(0.22),
             "spec-determinism Phase 2", size=8, color=MUTED)
    add_text(slide, Inches(12.0), Inches(7.1), Inches(0.5), Inches(0.22),
             str(number), size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def add_takeaway(slide, text: str) -> None:
    add_text(slide, Inches(1.55), Inches(6.55), Inches(10.25), Inches(0.34),
             text, size=13, bold=True, italic=True, color=TEAL,
             align=PP_ALIGN.CENTER)


def add_code_panel(slide, x, y, w, h, label: str, code: str, *, code_size=13):
    add_text(slide, x, y, w, Inches(0.38), label, size=17, bold=True, color=TEAL)

    # Offset rectangle for a subtle shadow, matching the reference deck.
    shadow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.04),
                                    y + Inches(0.68), w, h)
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = SHADOW
    shadow.line.fill.background()

    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.64), w, h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = CODE_BG
    panel.line.color.rgb = CODE_BORDER
    panel.line.width = Pt(0.7)

    box = slide.shapes.add_textbox(
        x + Inches(0.22), y + Inches(0.79), w - Inches(0.42), h - Inches(0.22)
    )
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.text = code
    for paragraph in tf.paragraphs:
        paragraph.space_after = Pt(0)
        for run in paragraph.runs:
            run.font.name = CODE_FONT
            run.font.size = Pt(code_size)
            run.font.color.rgb = DARK
    return panel


def add_bullets(slide, items: list[str], x, y, w, h, *, size=21):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    first = True
    for item in items:
        lines = item.split("\n")
        for idx, line in enumerate(lines):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(10) if idx == len(lines) - 1 else Pt(2)
            run = p.add_run()
            run.text = ("- " if idx == 0 else "  ") + line
            run.font.name = BODY_FONT
            run.font.size = Pt(size)
            run.font.color.rgb = DARK
    return box


def add_table(slide, x, y, w, h, headers: list[str], rows: list[list[str]],
              *, size=13):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h)
    table = shape.table
    for col in table.columns:
        col.width = int(w / len(headers))
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = header
        p.runs[0].font.name = BODY_FONT
        p.runs[0].font.size = Pt(size)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE
    for r, row in enumerate(rows, 1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PALE_TEAL if r % 2 else WHITE
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = value
            p.runs[0].font.name = BODY_FONT
            p.runs[0].font.size = Pt(size)
            p.runs[0].font.color.rgb = DARK
    return shape


def blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    return slide


def slide_title(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(0.02), SLIDE_W, Inches(0.025)
    )
    top.fill.solid()
    top.fill.fore_color.rgb = TEAL
    top.line.fill.background()
    bottom = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(7.18), SLIDE_W, Inches(0.025)
    )
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = TEAL
    bottom.line.fill.background()
    add_text(
        slide,
        Inches(0.65), Inches(2.15), Inches(12.0), Inches(0.75),
        "spec-determinism Phase 2",
        size=34, bold=True, color=TEAL, font=TITLE_FONT, align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        Inches(0.65), Inches(3.45), Inches(12.0), Inches(0.35),
        "View Registry -> PR-G",
        size=18, color=MUTED, align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        Inches(0.65), Inches(4.08), Inches(12.0), Inches(0.3),
        "Tianyu Chen · May 2026",
        size=14, color=MUTED, align=PP_ALIGN.CENTER,
    )


def slide_problem_definition(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    add_header(slide, "Problem Definition: Spec Determinism")
    add_code_panel(
        slide, Inches(0.72), Inches(1.25), Inches(5.75), Inches(2.55),
        "Question (one spec, two candidate outputs)",
        "Spec of fn f : X -> T contributes\n"
        "    P(x)     := requires clause\n"
        "    Q(x, y)  := ensures clause\n\n"
        "f is deterministic on x iff:\n"
        "    forall x, y1, y2.\n"
        "      P(x) && Q(x, y1) && Q(x, y2)\n"
        "        ==> equal_T(y1, y2)\n\n"
        "Same P, Q; only the output is re-bound.",
        code_size=12,
    )
    add_code_panel(
        slide, Inches(6.95), Inches(1.25), Inches(5.75), Inches(2.55),
        "Generated proof obligation",
        "proof fn check(x, y1, y2)\n"
        "  requires\n"
        "    P(x),\n"
        "    Q(x, y1),   // ensures on y1\n"
        "    Q(x, y2),   // ensures on y2\n"
        "  ensures\n"
        "    equal_T(y1, y2)\n"
        "{ }",
        code_size=12,
    )
    add_bullets(
        slide,
        [
            "There is only ONE ensures predicate Q; we re-bind its output to y1 and y2 to ask 'must they agree?'.",
            "equal_T is generated from the return type and equality policy, not from the spec text.",
            "Witness (counterexample to the implication) -> spec leaves output dimensions unconstrained.",
            "verus_error -> tool failed to express or check the obligation.",
        ],
        Inches(1.0), Inches(4.55), Inches(11.3), Inches(1.6), size=17,
    )
    add_takeaway(slide, "Determinism = Q is functional in y given P(x), up to type-directed spec equality equal_T.")
    add_footer(slide, n)


def slide_example_problem(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    add_header(slide, "Motivating Example: Runtime Equality != Spec Equivalence")
    add_code_panel(
        slide, Inches(0.62), Inches(1.18), Inches(5.55), Inches(2.92),
        "Spec shape",
        "pub struct AbstractEndPoint {\n"
        "  pub id: Vec<u8>\n"
        "}\n\n"
        "fn make_endpoint(bytes: Vec<u8>)\n"
        "  -> (r: AbstractEndPoint)\n"
        "ensures\n"
        "  r.id@ == bytes@\n",
        code_size=13,
    )
    add_code_panel(
        slide, Inches(6.88), Inches(1.18), Inches(5.55), Inches(2.92),
        "Naive determinism test",
        "spec fn equal(r1: AbstractEndPoint,\n"
        "              r2: AbstractEndPoint) -> bool {\n"
        "  r1 == r2\n"
        "}\n\n"
        "// compares Vec ptr / cap / len\n"
        "// z3 can report a spurious witness\n",
        code_size=13,
    )
    add_takeaway(slide, "The spec pins bytes@, but raw equality can still compare runtime representation noise.")
    add_footer(slide, n)


def slide_view_fix(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    add_header(slide, "A-2 Follow-up: equal_T Too Strong After A-1 (same function)")
    add_code_panel(
        slide, Inches(0.62), Inches(1.18), Inches(6.0), Inches(2.95),
        "Before A-2: ghost lifted, but compared element-wise",
        "// after A-1: range_alloc_and_map_io equal_T is well-formed\n"
        "equal((u1, g1), (u2, g2))\n"
        "  := u1 == u2 && (g1)@ == (g2)@\n"
        "// Seq<PagePtr> compared element-wise\n\n"
        "ensures only force:  self.wf()\n"
        "(allocator may pick any free page)\n\n"
        "z3 witness (classic A-2 signature):\n"
        "  r1 = (1, ghost(seq![PagePtr(0x1000)]))\n"
        "  r2 = (1, ghost(seq![PagePtr(0x2000)]))\n"
        "  // |r1@| == |r2@| == 1; r1@ != r2@\n"
        "result: ok_with_witness",
        code_size=10,
    )
    add_code_panel(
        slide, Inches(6.78), Inches(1.18), Inches(6.0), Inches(2.95),
        "After A-2: coarser det_view for allocator output",
        "// view registry: physical pages -> abstract\n"
        "// page-frame state; allocator choice is opaque.\n"
        "impl DetView for Ghost<Seq<PagePtr>> {\n"
        "  spec fn det_view(g: Ghost<Seq<PagePtr>>)\n"
        "    -> Multiset<PageSize> {\n"
        "    Multiset::from_seq((g)@.map(size_of))\n"
        "  }\n"
        "}\n\n"
        "equal((u1, g1), (u2, g2))\n"
        "  := u1 == u2 && det_view(g1) =~= det_view(g2)\n\n"
        "z3: UNSAT  (same size-class multiset)\n"
        "result: ok_without_witness",
        code_size=10,
    )
    add_bullets(
        slide,
        [
            "After A-1 unblocked Ghost handling, equal_T is well-formed but over-strong:\n"
            "it forces every PagePtr to match across runs, while ensures only pin self.wf().\n"
            "z3 reports the canonical A-2 signature -- same length, different elements.",
            "Two remediation paths for an A-2 witness:\n"
            "  (1) user-side: strengthen ensures so the values are pinned\n"
            "      (turns into A-1-real true positive);\n"
            "  (2) tool-side: register a coarser det_view in the view registry\n"
            "      that drops intentional allocator nondeterminism.",
            "Effect of path (2): the same function flips from ok_with_witness to\n"
            "ok_without_witness without rewriting any Verus proofs --\n"
            "the determinism graph just stops observing allocator-chosen page ids.",
        ],
        Inches(0.7), Inches(4.25), Inches(12.0), Inches(2.05), size=12,
    )
    add_takeaway(slide, "A-1 fixes the tool; A-2 aligns equal_T with the spec author's actual notion of 'same outcome'.")
    add_footer(slide, n)


def slide_status(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    add_header(slide, "Corpus Numbers: Baseline -> Now -> Remaining")
    add_code_panel(
        slide, Inches(0.7), Inches(1.18), Inches(12.0), Inches(0.95),
        "Buckets (status x assumes)",
        "ok_without_witness = status==ok && assumes==[]   (clean determinism proof; our target)\n"
        "ok_with_witness    = status==ok && assumes!=[]   (A-2 false positive: spec ok, equal_T over-specifies)\n"
        "verus_error        = generated obligation rejected by Verus       (A-1 tool failure)",
        code_size=10,
    )
    add_table(
        slide, Inches(0.7), Inches(2.25), Inches(12.0), Inches(1.35),
        ["Bucket", "Axis", "Baseline (42c1248, 04-29)", "Now (33bd09a, 05-11)", "Delta", "Still to fix"],
        [
            ["ok_without_witness", "-- (target)", "1079", "1090", "+11", "--"],
            ["ok_with_witness", "A-2", "376", "366", "-10", "366"],
            ["verus_error", "A-1", "191", "190", "-1", "190"],
            ["TOTAL", "-", "1647", "1647", "0", "556"],
        ],
        size=10,
    )
    add_table(
        slide, Inches(0.7), Inches(3.85), Inches(12.0), Inches(2.5),
        ["Project (baseline)", "n", "ok_with_witness (A-2)", "verus_error (A-1)", "Wins so far"],
        [
            ["atmosphere", "1363", "289", "100", "1 (PageMap)"],
            ["ironkv", "214", "76", "44", "1 (Constants)"],
            ["memory-allocator", "16", "9", "1", "8 (CommitMask cohort)"],
            ["nrkernel", "8", "1", "2", "1 (ArchExec)"],
            ["vest", "2", "1", "0", "0"],
            ["storage", "43", "0", "43", "0"],
            ["anvil-library", "1", "0", "1", "0"],
            ["TOTAL", "1647", "376", "191", "11"],
        ],
        size=10,
    )
    add_takeaway(slide, "11 wins locked in; 366 A-2 + 190 A-1 remain. Active view-registry rerun (4eb7376, atmosphere 944/1363) targets A-1 toward ~130.")
    add_footer(slide, n)


def slide_fp_nested_view(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    add_header(slide, "FP-B Example: Nested-Container View Uninterpretation")
    add_code_panel(
        slide, Inches(0.62), Inches(1.18), Inches(6.05), Inches(3.05),
        "ironkv::clone_option_vec_u8",
        "#[verifier::external_body]\n"
        "pub fn clone_option_vec_u8(\n"
        "    ov: Option<&Vec<u8>>,\n"
        ") -> (res: Option<Vec<u8>>)\n"
        "  ensures match ov {\n"
        "    Some(e1) => res.is_some()\n"
        "      && e1@ == res.get_Some_0()@,\n"
        "    None => res.is_None(),\n"
        "  }\n\n"
        "// codegen emits structural Vec eq:\n"
        "spec fn equal(r1, r2) -> bool {\n"
        "  ((r1 is Some) == (r2 is Some))\n"
        "  && ((r1 is Some) ==>\n"
        "        r1->Some_0 == r2->Some_0)\n"
        "}",
        code_size=11,
    )
    add_code_panel(
        slide, Inches(6.78), Inches(1.18), Inches(5.95), Inches(3.05),
        "Refined witness (7 assumes, post #14c)",
        "assume(ov is Some);\n"
        "assume(ov->Some_0@.len() == 0);\n"
        "assume(r1 is Some); assume(r2 is Some);\n"
        "assume(r1->Some_0@.len() == 0);\n"
        "assume(r2->Some_0@.len() == 0);\n"
        "assume(!equal(r1, r2));\n\n"
        "// (A) structural Option<Vec> equal:\n"
        "//    r1->Some_0 == r2->Some_0 ?\n"
        "//    --> z3 SAT  (spurious witness)\n\n"
        "// (B) recurse: instantiate Option<T>\n"
        "//     to inner Vec, compare views:\n"
        "//    r1->Some_0@ == r2->Some_0@ ?\n"
        "//    --> z3 UNSAT  (FP gone)",
        code_size=11,
    )
    add_bullets(
        slide,
        [
            "ensures forces r1@ == ov@ == r2@; the only freedom is the uninterpreted Vec eq inside Some.",
            "Z3 limitation: it does not peel Option<Vec<u8>> to its inner Vec on its own; lemma_vec_obeys_view_eq is not in the default broadcast group.",
            "Codegen does the peeling: recurse type_args, compare r1->Some_0@ instead of r1->Some_0. Same query -> UNSAT, FP gone.",
        ],
        Inches(0.85), Inches(4.4), Inches(11.95), Inches(2.0), size=14,
    )
    add_takeaway(slide, "FP-B is a z3 type-reasoning gap; codegen has to peel Option<Vec> and compare views itself.")
    add_footer(slide, n)


def slide_axes(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    add_header(slide, "A1-A3: Three Problem Classes")
    add_table(
        slide, Inches(0.85), Inches(1.25), Inches(11.65), Inches(2.2),
        ["Axis", "What goes wrong", "Failure shape"],
        [
            ["A-1", "type has no View, or its View is missing fields", "verus_error / partial witness"],
            ["A-2", "equal-fn compares too much (runtime structure)", "spurious ok_with_witness"],
            ["A-3", "equal-fn uses wrong semantics for the variant", "nested Result / Err mis-compare"],
        ],
        size=13,
    )
    add_bullets(
        slide,
        [
            "Three axes correspond to the three things that can go wrong when we turn determinism into a verifiable obligation.",
            "A-1 is upstream: without a view we cannot even ask the question -- z3 sees opaque types and verification stalls.",
            "A-2 is the canonical false-positive shape: the spec is determinism-equivalent, but equal_T over-specifies, so z3 finds a wiggle.",
            "A-3 is structural: the equal-fn picks the wrong relation for a sum / nested type (e.g., comparing two Err payloads byte-for-byte when only the tag matters).",
            "The following slides walk one concrete example per axis (A-1: missing ghost view; A-2: runtime noise; A-3: nested Err).",
        ],
        Inches(0.85), Inches(3.6), Inches(11.95), Inches(2.7), size=14,
    )
    add_takeaway(slide, "Every witness/error we see falls into exactly one of A-1, A-2, A-3 -- that triage drives the fix.")
    add_footer(slide, n)


def slide_solution_graph(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    add_header(slide, "Core Solution: View-Based Equal-Fn + Type Dependency Graph")
    add_bullets(
        slide,
        [
            "equal_T should compare spec views, not runtime structure -- runtime detail leaks implementation noise into determinism.",
            "atmosphere: every op points at a PageTable; we want equality modulo PageTable.view(), ignoring bit-allocation / arena layout.",
            "Per-function hand-written views do not scale -- need a repo-wide answer.",
            "Build a type dependency graph across the repo, attach a View to every node -> equal_T = recursive walk of that graph.",
        ],
        Inches(0.7), Inches(1.18), Inches(12.0), Inches(1.95), size=13,
    )
    add_code_panel(
        slide, Inches(0.65), Inches(3.4), Inches(6.05), Inches(2.55),
        "Type dependency graph (repo-wide)",
        "node : Verus type T\n"
        "edge : T contains / refers to U\n\n"
        "PageMapHost\n"
        "  -> tables: Map<NodeId, PageTable>\n"
        "        -> entries: Vec<PageEntry>\n"
        "              -> bits: u64    // runtime\n"
        "              -> mapped: Seq<Frame>  // spec\n\n"
        "View attached at every node\n"
        "  -> bit-level diffs collapse;\n"
        "     mapped Seq drives equality.",
        code_size=11,
    )
    add_code_panel(
        slide, Inches(6.95), Inches(3.4), Inches(5.95), Inches(2.55),
        "How the graph turns into equal_T",
        "equal_T(r1: T, r2: T) :=\n"
        "  let v1 = T.view(r1);\n"
        "  let v2 = T.view(r2);\n"
        "  forall field f of T::V.\n"
        "    equal_{type_of(f)}(v1.f, v2.f)\n\n"
        "Recurse on each child type using\n"
        "the same rule from the graph.\n"
        "Leaf prims fall back to ==.\n"
        "Missing view -> synthesize or flag.",
        code_size=11,
    )
    add_takeaway(slide, "Per-function equal_T comes for free once every type in the repo has a view in the dependency graph.")
    add_footer(slide, n)


def slide_solution_equalfn(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    add_header(slide, "Core Translation: Equal-Fn Becomes View Consistency")
    add_code_panel(
        slide, Inches(0.65), Inches(1.18), Inches(5.75), Inches(3.0),
        "Before",
        "spec fn equal(r1: T, r2: T) -> bool {\n"
        "  structural_eq(r1, r2)\n"
        "}\n\n"
        "// can compare private fields,\n"
        "// pointers, capacity,\n"
        "// wrapper internals,\n"
        "// or opaque Err payloads",
        code_size=13,
    )
    add_code_panel(
        slide, Inches(6.9), Inches(1.18), Inches(5.75), Inches(3.0),
        "After",
        "spec fn equal(r1: T, r2: T) -> bool {\n"
        "  view_eq_T(r1, r2)\n"
        "}\n\n"
        "view_eq_T(x, y):\n"
        "  if T has View: x.view() == y.view()\n"
        "  if container: forall element view_eq\n"
        "  if wrapper: unwrap, then recurse",
        code_size=12,
    )
    add_takeaway(slide, "The function-level check is generated mechanically from the return type's view graph.")
    add_footer(slide, n)


def slide_sources(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    add_header(slide, "Implementation: Build The View Registry")
    add_bullets(
        slide,
        [
            "The view registry is the concrete realization of the type dependency graph: a Type -> View mapping populated in four layers.",
            "Resolution is layered: L1 is consulted first, then L2 alias expansion, then L3 project-defined impls, and L4 LLM-synthesized views fill remaining holes.",
            "Each layer has a different trust level -- L1/L3 are accepted as-is, L2 is purely mechanical, L4 must pass the guardrail pipeline before being merged.",
        ],
        Inches(0.7), Inches(1.18), Inches(12.0), Inches(1.7), size=14,
    )
    add_table(
        slide, Inches(0.85), Inches(3.05), Inches(11.65), Inches(2.8),
        ["Layer", "Source", "Example", "Trust"],
        [
            ["L1", "hand-coded prelude views", "Vec -> Seq, HashMap -> Map", "trusted"],
            ["L2", "type alias expansion", "Pcid = usize -> reuse u64 view", "trusted"],
            ["L3", "project impl View for T", "atmosphere PageTable -> PageTableView", "trusted"],
            ["L4", "LLM-synthesized views", "structs with no existing view", "must pass guardrails"],
        ],
        size=12,
    )
    add_takeaway(slide, "Equal-fn generation only runs after every type reachable from the return value has a resolved view in this registry.")
    add_footer(slide, n)


def slide_guardrails(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    add_header(slide, "L4 Guardrails: How We Accept LLM-Synthesized Views")
    add_bullets(
        slide,
        [
            "Only L4 (LLM-synthesized) needs guardrails -- L1/L2/L3 are trusted by construction.",
            "Pipeline: synth -> static lints M1-M4 -> critic pass -> cache (accept) / quarantine (suspect) / reject log -> corpus rerun -> diff baseline vs candidate.",
            "Rejected and quarantined views are kept as durable failure data (_rejected.jsonl), so a future lint can retroactively scan past candidates.",
        ],
        Inches(0.7), Inches(1.18), Inches(12.0), Inches(1.55), size=14,
    )
    add_code_panel(
        slide, Inches(0.65), Inches(2.85), Inches(6.05), Inches(3.1),
        "Pipeline stages",
        "LLM synthesizes View candidate\n"
        "  -> static lints M1-M4\n"
        "       (cheap syntactic checks)\n"
        "  -> critic pass\n"
        "       (semantic / consistency)\n"
        "  -> cache OR quarantine OR reject\n"
        "  -> corpus rerun on the candidate\n"
        "  -> compare witnesses vs baseline\n"
        "  -> promote, hold, or roll back",
        code_size=12,
    )
    add_code_panel(
        slide, Inches(6.95), Inches(2.85), Inches(5.95), Inches(3.1),
        "Static lints M1-M4 (priority M3 > M2 > M4 > M1)",
        "M3  view body reads fields of\n"
        "    external_body / opaque types\n"
        "M2  extra @ piped through Ghost\n"
        "    into non-viewable heads\n"
        "M4  self-recursive V uses TView,\n"
        "    but body emits bare self.f@\n"
        "M1  field@ used on types that\n"
        "    have no registered View",
        code_size=12,
    )
    add_takeaway(slide, "Generation is cheap; LLM views only enter the registry after passing lints, critic, and corpus diff.")
    add_footer(slide, n)


def slide_lints(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    add_header(slide, "PR-D5 / PR-E: Static Lints")
    add_table(
        slide, Inches(0.78), Inches(1.25), Inches(11.8), Inches(3.05),
        ["Rule", "Rejects"],
        [
            ["M1", "field@ on types without registered View"],
            ["M2", "extra @ through Ghost into non-viewable heads"],
            ["M3", "view body reads fields of external_body / opaque types"],
            ["M4", "self-recursive V uses TView but body emits bare self.f@"],
        ],
        size=13,
    )
    add_code_panel(
        slide, Inches(2.3), Inches(4.8), Inches(8.5), Inches(0.6),
        "Lint priority",
        "M3 > M2 > M4 > M1",
        code_size=18,
    )
    add_footer(slide, n)


def slide_recursive_bug(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    add_header(slide, "Example: Self-Recursive View Bug (caught by M4)")
    add_code_panel(
        slide, Inches(0.62), Inches(1.18), Inches(6.0), Inches(2.95),
        "Problem: bare @ does not descend",
        "pub struct PTDirView {\n"
        "  entries: Seq<Option<PTDirView>>\n"
        "}\n\n"
        "impl View for PTDir {\n"
        "  type V = PTDirView;\n"
        "  fn view(&self) -> PTDirView {\n"
        "    PTDirView { entries: self.entries@ }\n"
        "    //                          ^^^^^^^^^^\n"
        "    // type is Seq<Option<PTDir>>,\n"
        "    // not Seq<Option<PTDirView>>\n"
        "  }\n"
        "}",
        code_size=11,
    )
    add_code_panel(
        slide, Inches(6.78), Inches(1.18), Inches(6.0), Inches(2.95),
        "Fix: recurse view explicitly",
        "fn view(&self) -> PTDirView {\n"
        "  PTDirView {\n"
        "    entries: self.entries@.map(\n"
        "      |opt: Option<PTDir>|\n"
        "        opt.map(|d| d.view())\n"
        "    )\n"
        "  }\n"
        "}\n\n"
        "// now Seq<Option<PTDirView>>\n"
        "// matches the declared field type",
        code_size=11,
    )
    add_bullets(
        slide,
        [
            "Why a bug: `self.entries@` only lifts the outer Vec to Seq;\n"
            "the inner PTDir stays in runtime form, so equal_T silently\n"
            "compares runtime PTDir instead of PTDirView (abstraction lost, no compile error).",
            "Detection: M4 compares the declared V field type (Seq<Option<PTDirView>>)\n"
            "against the type that `self.f@` actually produces (Seq<Option<PTDir>>);\n"
            "mismatch -> reject the view.",
            "Fix: rewrite the body to map through the recursive slot,\n"
            "calling .view() on the inner PTDir so types line up.",
        ],
        Inches(0.7), Inches(4.25), Inches(12.0), Inches(2.05), size=12,
    )
    add_takeaway(slide, "M4 turns a silent type-mismatch (works at compile time, wrong at spec time) into a hard lint failure.")
    add_footer(slide, n)


def slide_prf(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    add_header(slide, "A-1 Example: Missing Ghost View (atmosphere::range_alloc_and_map_io)")
    add_code_panel(
        slide, Inches(0.62), Inches(1.18), Inches(6.0), Inches(2.95),
        "Before: Ghost not in registry -> verus_error",
        "// atmosphere/.../range_alloc_and_map_io.rs:2222\n"
        "pub fn range_alloc_and_map_io(&mut self,\n"
        "    target: ProcPtr, va_range: &VaRange4K,\n"
        ") -> (ret: (usize, Ghost<Seq<PagePtr>>))\n"
        "  requires old(self).wf(), ...\n"
        "  ensures  self.wf()  // other clauses commented out\n\n"
        "// codegen:\n"
        "equal(r1, r2) := r1.0 == r2.0 && r1.1 == r2.1\n"
        "// Ghost -> UNKNOWN fallback -> structural ==\n"
        "// Verus refused: 'no equality on Ghost'\n"
        "result: verus_error",
        code_size=10,
    )
    add_code_panel(
        slide, Inches(6.78), Inches(1.18), Inches(6.0), Inches(2.95),
        "After: register Ghost<T> -> T view -> witness shows up",
        "equal((u1, g1), (u2, g2)) :=\n"
        "  u1 == u2 && (g1)@ == (g2)@\n"
        "// Ghost branch strips wrapper -> compares Seq<PagePtr>\n\n"
        "z3 witness:\n"
        "  r1 = (1, ghost(seq![PagePtr(0x1000)]))\n"
        "  r2 = (1, ghost(seq![PagePtr(0x2000)]))\n"
        "// both satisfy self.wf();\n"
        "// ensures says nothing about which pages\n\n"
        "result: ok_with_witness\n"
        "  (real A-2 gap surfaced)",
        code_size=10,
    )
    add_bullets(
        slide,
        [
            "Problem (A-1): Ghost<T> had no entry in the view registry,\n"
            "so equal_T fell back to structural Ghost == Ghost,\n"
            "which Verus rejects -> verus_error hides the determinism question.",
            "Fix: add a prelude rule Ghost<T> -> T.view();\n"
            "equal_T now compares the inner Seq<PagePtr> view.",
            "Effect: tool stops crashing and exposes a real gap --\n"
            "ensures pins only self.wf(), nothing about which page sequence\n"
            "was allocated, so z3 finds two distinct PagePtr seqs.",
        ],
        Inches(0.7), Inches(4.25), Inches(12.0), Inches(2.05), size=12,
    )
    add_takeaway(slide, "Closing the Ghost-view gap turns a tool failure (A-1) into a legitimate spec-incompleteness witness (A-2).")
    add_footer(slide, n)


def slide_prg(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    add_header(slide, "A-3 Example: Nested Err Semantics (batch_lookup)")
    add_code_panel(
        slide, Inches(0.62), Inches(1.18), Inches(6.0), Inches(2.95),
        "Before: structural Seq == Seq compares Err payloads",
        "fn batch_lookup(...) -> (r: Seq<Result<u32, MyErr>>)\n"
        "  policy: errs_equivalent = true\n\n"
        "// codegen emitted plain r1 == r2;\n"
        "// SEQ branch fell through to structural ==,\n"
        "// so Err payloads got compared byte-for-byte.\n\n"
        "z3 witness (spurious):\n"
        "  r1 = seq![Err(timeout), Err(parse)]\n"
        "  r2 = seq![Err(parse),   Err(timeout)]\n"
        "  // policy says these are equivalent,\n"
        "  // but r1 != r2 structurally -> witness",
        code_size=10,
    )
    add_code_panel(
        slide, Inches(6.78), Inches(1.18), Inches(6.0), Inches(2.95),
        "After: elementwise Result eq + policy collapse",
        "equal(r1, r2) :=\n"
        "  r1.len() == r2.len()\n"
        "  && forall |i: int| 0 <= i < r1.len() ==>\n"
        "       match (r1[i], r2[i]) {\n"
        "         (Ok(a), Ok(b))     => a == b,\n"
        "         (Err(_), Err(_))   => true,  // policy\n"
        "         _                  => false,\n"
        "       }\n\n"
        "on r1, r2 above: all i match (Err, Err) -> true\n"
        "z3: UNSAT  ->  result: ok_without_witness",
        code_size=10,
    )
    add_bullets(
        slide,
        [
            "Problem (A-3): TypeKind.SEQ was in the primitive-`==` list,\n"
            "so Seq<Result<U, Err>> compared Err byte-for-byte\n"
            "even though policy declared errs_equivalent = true.",
            "Fix: new helpers _contains_result / _container_needs_elementwise;\n"
            "explicit SEQ branch emits len== + forall element-eq when element\n"
            "contains Result and the policy collapses Err. Same fix for Map values.",
            "Effect: Err-payload noise no longer blocks the tool;\n"
            "Ok-side incompleteness becomes visible -- e.g., r1=[Ok(1)] vs r2=[Ok(2)] still SAT (real A-2).",
        ],
        Inches(0.7), Inches(4.25), Inches(12.0), Inches(2.05), size=12,
    )
    add_takeaway(slide, "A-3 fixed by recursing equality through Seq/Map containers and honouring the policy at the leaf Result node.")
    add_footer(slide, n)


def slide_numbers(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    add_header(slide, "Numbers")
    add_code_panel(
        slide, Inches(0.78), Inches(1.2), Inches(5.55), Inches(2.0),
        "Baseline: 42c1248, 2026-04-29",
        "n=1647\n"
        "ok=1455\n"
        "witness=376\n"
        "verus_error=191\n"
        "runner_crash=1",
        code_size=16,
    )
    add_code_panel(
        slide, Inches(6.95), Inches(1.2), Inches(5.55), Inches(2.0),
        "Post D5/E: 33bd09a, 2026-05-11",
        "ok=1456 (+1)\n"
        "witness=366 (-10)\n"
        "verus_error=190 (-1)\n\n"
        "11 true wins\n"
        "0 regressions",
        code_size=16,
    )
    add_takeaway(slide, "witness means ok_with_witness, a subset of raw ok.")
    add_footer(slide, n)


def slide_current(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    add_header(slide, "Current Wrapper / Nested-Err Rerun")
    add_table(
        slide, Inches(1.0), Inches(1.35), Inches(11.0), Inches(2.35),
        ["Item", "Value"],
        [
            ["commit", "4eb7376"],
            ["atmosphere progress", "944 / 1363 targets, 69%"],
            ["expected verus_error movement", "toward the ~130 region"],
            ["per-target cost", "4.89s -> 7.93s (+62%)"],
        ],
        size=15,
    )
    add_takeaway(slide, "The extra cost is the expected tradeoff for richer schemas and finer narrowing.")
    add_footer(slide, n)


def slide_takeaways(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    add_header(slide, "Takeaways: What We Did")
    add_bullets(
        slide,
        [
            "Started by debugging spec-determinism on a 1647-target verusage corpus;\n"
            "raw output gave only ok / verus_error and was noisy.",
            "Triaged every failure into three problem classes:\n"
            "A-1 missing/wrong view (tool fails),\n"
            "A-2 equal-fn too strict (spurious witness),\n"
            "A-3 wrong equality semantics (Err / nested types).",
            "Proposed a unified fix: equal_T compares views,\n"
            "and views come from a repo-wide type dependency graph\n"
            "populated by L1 prelude / L2 alias / L3 project / L4 LLM-synthesized layers.",
            "Hardened L4 with a guardrail pipeline (lints M1-M4 + critic + quarantine + corpus diff)\n"
            "so synthesized views never silently break the corpus.",
            "Result: 11 true wins, 0 regressions, every failure now lands in a named class\n"
            "with a known fix path -- the rest is execution.",
        ],
        Inches(0.8), Inches(1.25), Inches(12.0), Inches(5.2), size=15,
    )
    add_footer(slide, n)


def slide_next(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    add_header(slide, "Next: Drive ok_with_witness and verus_error to Zero")
    add_bullets(
        slide,
        [
            "Core plan: walk every remaining ok_with_witness and verus_error,\n"
            "until both columns are exhausted.",
            "366 ok_with_witness (A-2): for each, run the FP-B litmus (z3 probe with view-aligned narrows).\n"
            "  - UNSAT under view-aligned narrows -> codegen FP, fix by view-lifting equal_T.\n"
            "  - SAT survives -> real spec incompleteness, surface to spec author.",
            "190 verus_error (A-1): classify each obligation by failure shape and clear the registry hole.\n"
            "  - missing wrapper view (Ghost/Tracked/PointsTo) -> add prelude rule.\n"
            "  - opaque / external_body type -> add a view or skip with a logged reason.\n"
            "  - newtype-of-usize / type alias -> finish L2 unwrap.",
            "Supporting work:\n"
            "  - implement FP-B codegen fix (recursive view-lift through nested type_args).\n"
            "  - finish the in-flight view-registry rerun (atmosphere 944/1363) and lock in the new baseline.\n"
            "  - retry the quarantined L4 views after M1-M4 + critic are strict enough.",
        ],
        Inches(0.8), Inches(1.25), Inches(12.0), Inches(5.2), size=14,
    )
    add_footer(slide, n)


def slide_true_positive(prs: Presentation, n: int) -> None:
    slide = blank(prs)
    add_header(slide, "True Positive: Under-Constrained Ensures (CommitMask::next_run)")
    add_code_panel(
        slide, Inches(0.62), Inches(1.18), Inches(6.0), Inches(2.95),
        "Spec: only pins 'the run is committed', not which run",
        "pub fn next_run(&self, idx: usize)\n"
        "  -> (res: (usize, usize))\n"
        "  requires 0 <= idx < COMMIT_MASK_BITS,\n"
        "  ensures ({ let (next_idx, count) = res;\n"
        "    next_idx + count <= COMMIT_MASK_BITS\n"
        "    && forall |t| next_idx <= t < next_idx + count\n"
        "         ==> self@.contains(t)\n"
        "  })\n\n"
        "// missing: count must be MAXIMAL\n"
        "// missing: next_idx must be the FIRST run >= idx\n"
        "// missing: lower bound count > 0 unless no run exists",
        code_size=10,
    )
    add_code_panel(
        slide, Inches(6.78), Inches(1.18), Inches(6.0), Inches(2.95),
        "Witness: two outputs satisfy the same ensures",
        "assume self@ = {0, 1, 2, 3}, idx = 0\n\n"
        "candidate r1 = (0, 0)\n"
        "  forall |t| 0 <= t < 0 -> contains(t)\n"
        "  vacuously true   -> ensures holds\n\n"
        "candidate r2 = (0, 4)\n"
        "  forall |t| 0 <= t < 4 -> contains(t)\n"
        "  {0,1,2,3} all present -> ensures holds\n\n"
        "ensures(r1) /\\ ensures(r2) /\\ r1 != r2\n"
        "  -> determinism fails, real witness",
        code_size=10,
    )
    add_bullets(
        slide,
        [
            "Tool stack is correct: equal_T = tuple struct ==,\n"
            "view registry already has CommitMask -> Set<int>;\n"
            "the gap is in the spec -- it never says which run is returned.",
            "Source comments (commit_mask.rs:88-90) already flag the missing clauses;\n"
            "the tool surfaced a latent ambiguity, not a tool bug.",
            "Fix path: hand the witness back to the spec author.\n"
            "Add 'count is maximal' + 'no committed bit in [idx, next_idx)'.",
        ],
        Inches(0.7), Inches(4.25), Inches(12.0), Inches(2.05), size=12,
    )
    add_takeaway(slide, "True positive: a real determinism gap that the spec author must close -- not a tool or codegen artefact.")
    add_footer(slide, n)


def build_deck(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        slide_problem_definition,
        slide_axes,
        slide_fp_nested_view,
        slide_solution_graph,
        slide_sources,
        slide_guardrails,
        slide_recursive_bug,
        slide_prf,
        slide_view_fix,
        slide_prg,
        slide_true_positive,
        slide_status,
        slide_next,
    ]
    for i, builder in enumerate(builders, 1):
        builder(prs, i)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> int:
    default_output = Path(__file__).with_name("spec-determinism-phase2-2026-05-12.pptx")
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, default=default_output,
                        help=f"Output pptx path (default: {default_output})")
    args = parser.parse_args()
    build_deck(args.output)
    print(args.output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
